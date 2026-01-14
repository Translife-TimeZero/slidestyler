"""
World-Class PPTX Exporter
Creates stunning, professional PowerPoint presentations with modern design
"""

import os
from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap, qn
import re


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    return RGBColor(0, 0, 0)


def lighten_color(hex_color: str, factor: float = 0.3) -> str:
    """Lighten a hex color"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02x}{g:02x}{b:02x}"


def darken_color(hex_color: str, factor: float = 0.2) -> str:
    """Darken a hex color"""
    hex_color = hex_color.lstrip('#')
    r = int(int(hex_color[0:2], 16) * (1 - factor))
    g = int(int(hex_color[2:4], 16) * (1 - factor))
    b = int(int(hex_color[4:6], 16) * (1 - factor))
    return f"#{r:02x}{g:02x}{b:02x}"


class WorldClassExporter:
    """Creates stunning, world-class PowerPoint presentations"""
    
    def __init__(self, style: Dict, slides_data: List[Dict]):
        self.style = style
        self.theme = style.get('theme', {})
        self.typography = style.get('typography', {})
        self.slides_data = slides_data
        self.slide_count = len(slides_data)
        
        # Enhanced color palette
        self.primary = self.theme.get('primary', '#0077b6')
        self.accent = self.theme.get('accent', '#00b4d8')
        self.background = self.theme.get('background', '#ffffff')
        self.text_color = self.theme.get('text', '#1a1a2e')
        self.text_muted = self.theme.get('text_muted', '#64748b')
        self.surface = self.theme.get('surface', '#f8fafc')
        
        # Design parameters
        self.title_font = self.typography.get('heading', 'Segoe UI')
        self.body_font = self.typography.get('body', 'Segoe UI')
        
        # Create presentation with 16:9 aspect ratio
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def export(self, output_path: str) -> str:
        """Export all slides to PPTX with world-class design"""
        for i, slide_data in enumerate(self.slides_data):
            self._create_world_class_slide(slide_data, i)
        
        self.prs.save(output_path)
        return output_path
    
    def _create_world_class_slide(self, slide_data: Dict, index: int):
        """Create a stunning slide based on content and position"""
        layout_type = slide_data.get('layout_type', 'content')
        original_content = slide_data.get('original_content', [])
        
        # Extract and clean content
        title, subtitle, body_texts = self._extract_content(original_content)
        
        # First slide is always a hero title
        if index == 0:
            self._create_hero_title_slide(title, subtitle, body_texts)
        # Last slide is a closing slide
        elif index == self.slide_count - 1 and self.slide_count > 2:
            self._create_stunning_closing_slide(title, body_texts)
        # Determine best layout based on content
        elif len(body_texts) >= 6:
            self._create_grid_content_slide(title, body_texts)
        elif len(body_texts) >= 4:
            self._create_two_column_modern_slide(title, body_texts)
        elif any(self._looks_like_stat(t) for t in body_texts):
            self._create_stats_showcase_slide(title, body_texts)
        else:
            self._create_elegant_content_slide(title, body_texts)
    
    def _extract_content(self, original_content: List) -> Tuple[str, str, List[str]]:
        """Extract and clean content from slide data"""
        title = ""
        subtitle = ""
        body_texts = []
        
        skip_types = ['sldNum', 'ftr', 'dt', 'hdr']
        title_types = ['title', 'ctrTitle', 'TITLE', 'CENTER_TITLE']
        subtitle_types = ['subTitle', 'SUBTITLE', 'subtitle']
        
        for item in original_content:
            if isinstance(item, str):
                text = item.strip()
                item_type = 'body'
            else:
                item_type = item.get('type', 'body')
                text = item.get('text', '').strip()
            
            if not text or item_type in skip_types:
                continue
            if text.isdigit() and len(text) <= 3:
                continue
                
            if item_type in title_types and not title:
                title = text
            elif item_type in subtitle_types:
                subtitle = text
            else:
                body_texts.append(text)
        
        # Use first body text as title if none found
        if not title and body_texts:
            title = body_texts.pop(0)
        
        return title, subtitle, body_texts
    
    def _looks_like_stat(self, text: str) -> bool:
        """Check if text looks like a statistic"""
        return bool(re.search(r'\d+[%$€£]|\d+\s*(percent|million|billion|k\b|m\b)', text.lower()))
    
    # ==================== HERO TITLE SLIDE ====================
    
    def _create_hero_title_slide(self, title: str, subtitle: str, body_texts: List[str]):
        """Create a stunning hero title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Dark gradient-like background using primary color
        self._set_slide_background(slide, self.primary)
        
        # Add geometric accent shapes
        self._add_hero_decorations(slide)
        
        # Large accent circle (top right, partially visible)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(10), Inches(-2), 
            Inches(5), Inches(5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.primary, 0.15))
        circle.line.fill.background()
        
        # Small accent circle
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(-1), Inches(5), 
            Inches(3), Inches(3)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = hex_to_rgb(self.accent)
        circle2.line.fill.background()
        
        # Title - large and bold
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(10), Inches(2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.title_font
        
        # Subtitle or first body text
        sub_text = subtitle or (body_texts[0] if body_texts else "")
        if sub_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(8), Inches(1))
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sub_text
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.body_font
        
        # Decorative line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.8), Inches(4.1), 
            Inches(2), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.accent)
        line.line.fill.background()
    
    def _add_hero_decorations(self, slide):
        """Add subtle geometric decorations to hero slide"""
        # Top-right corner accent
        corner = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(11.5), Inches(0),
            Inches(1.833), Inches(1.5)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = hex_to_rgb(self.accent)
        corner.line.fill.background()
        corner.rotation = 90
    
    # ==================== ELEGANT CONTENT SLIDE ====================
    
    def _create_elegant_content_slide(self, title: str, body_texts: List[str]):
        """Create an elegant content slide with modern design"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, self.background)
        
        # Left accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.15), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(self.primary)
        accent_bar.line.fill.background()
        
        # Section indicator circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.4), Inches(0.5),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(self.accent)
        circle.line.fill.background()
        
        # Title with accent underline
        if title:
            title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(1.2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # Clean up title (remove extra parts after common delimiters)
            clean_title = title.split('.')[0].split('BAGAIMANA')[0].strip()
            if len(clean_title) > 80:
                clean_title = clean_title[:77] + "..."
            p.text = clean_title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.primary)
            p.font.name = self.title_font
            
            # Underline accent
            underline = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.2), Inches(1.7),
                Inches(1.5), Inches(0.05)
            )
            underline.fill.solid()
            underline.fill.fore_color.rgb = hex_to_rgb(self.accent)
            underline.line.fill.background()
        
        # Content cards
        if body_texts:
            self._add_content_cards(slide, body_texts[:6], start_y=2.0)
    
    def _add_content_cards(self, slide, items: List[str], start_y: float = 2.0):
        """Add content as elegant cards"""
        card_height = 0.75
        spacing = 0.15
        
        for i, text in enumerate(items):
            y = start_y + i * (card_height + spacing)
            if y > 6.5:
                break
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.2), Inches(y),
                Inches(11), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(self.surface)
            card.line.fill.background()
            
            # Accent indicator
            indicator = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.2), Inches(y),
                Inches(0.08), Inches(card_height)
            )
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = hex_to_rgb(self.accent)
            indicator.line.fill.background()
            
            # Card text
            text_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y + 0.15),
                Inches(10.5), Inches(card_height - 0.3)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # Truncate long text
            display_text = text if len(text) < 120 else text[:117] + "..."
            p.text = display_text
            p.font.size = Pt(15)
            p.font.color.rgb = hex_to_rgb(self.text_color)
            p.font.name = self.body_font
    
    # ==================== TWO COLUMN MODERN SLIDE ====================
    
    def _create_two_column_modern_slide(self, title: str, body_texts: List[str]):
        """Create a modern two-column layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, self.background)
        
        # Top accent bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.1)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = hex_to_rgb(self.primary)
        top_bar.line.fill.background()
        
        # Title
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            clean_title = title.split('.')[0].strip()[:70]
            p.text = clean_title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.primary)
            p.font.name = self.title_font
        
        # Split content into two columns
        mid = len(body_texts) // 2
        left_items = body_texts[:mid] if mid > 0 else body_texts[:2]
        right_items = body_texts[mid:] if mid > 0 else body_texts[2:]
        
        # Left column header
        left_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.7),
            Inches(5.5), Inches(0.5)
        )
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = hex_to_rgb(self.primary)
        left_header.line.fill.background()
        
        # Left content
        self._add_column_items(slide, left_items[:4], x=0.8, start_y=2.4)
        
        # Right column header
        right_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.9), Inches(1.7),
            Inches(5.5), Inches(0.5)
        )
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = hex_to_rgb(self.accent)
        right_header.line.fill.background()
        
        # Right content
        self._add_column_items(slide, right_items[:4], x=6.9, start_y=2.4)
    
    def _add_column_items(self, slide, items: List[str], x: float, start_y: float):
        """Add items to a column with icons"""
        for i, text in enumerate(items):
            y = start_y + i * 1.2
            
            # Icon circle
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(0.4), Inches(0.4)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.primary, 0.7))
            icon.line.fill.background()
            
            # Number in icon
            num_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.05), Inches(0.4), Inches(0.35))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.primary)
            p.alignment = PP_ALIGN.CENTER
            
            # Text
            text_box = slide.shapes.add_textbox(Inches(x + 0.55), Inches(y), Inches(4.8), Inches(1))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text[:100] if len(text) > 100 else text
            p.font.size = Pt(14)
            p.font.color.rgb = hex_to_rgb(self.text_color)
            p.font.name = self.body_font
    
    # ==================== GRID CONTENT SLIDE ====================
    
    def _create_grid_content_slide(self, title: str, body_texts: List[str]):
        """Create a grid layout for many items"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, self.background)
        
        # Side accent
        side_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(12.833), Inches(0),
            Inches(0.5), Inches(7.5)
        )
        side_accent.fill.solid()
        side_accent.fill.fore_color.rgb = hex_to_rgb(self.primary)
        side_accent.line.fill.background()
        
        # Title
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.5), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title.split('.')[0].strip()[:60]
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.primary)
            p.font.name = self.title_font
        
        # Create 2x3 or 3x2 grid
        items = body_texts[:6]
        cols = 3
        rows = 2
        card_width = 3.8
        card_height = 2.2
        start_x = 0.6
        start_y = 1.5
        gap = 0.3
        
        for i, text in enumerate(items):
            col = i % cols
            row = i // cols
            
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + gap)
            
            # Card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(card_width), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(self.surface)
            card.line.fill.background()
            
            # Top accent on card
            card_accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(card_width), Inches(0.08)
            )
            card_accent.fill.solid()
            # Alternate colors
            accent_colors = [self.primary, self.accent, darken_color(self.primary, 0.1)]
            card_accent.fill.fore_color.rgb = hex_to_rgb(accent_colors[i % 3])
            card_accent.line.fill.background()
            
            # Card number
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.2), Inches(y + 0.3),
                Inches(0.45), Inches(0.45)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = hex_to_rgb(accent_colors[i % 3])
            num_circle.line.fill.background()
            
            num_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.35), Inches(0.45), Inches(0.4))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Card text
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.9),
                Inches(card_width - 0.3), Inches(card_height - 1.1)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text[:90] if len(text) > 90 else text
            p.font.size = Pt(12)
            p.font.color.rgb = hex_to_rgb(self.text_color)
            p.font.name = self.body_font
    
    # ==================== STATS SHOWCASE SLIDE ====================
    
    def _create_stats_showcase_slide(self, title: str, body_texts: List[str]):
        """Create a visually striking statistics slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, self.primary)
        
        # Decorative circles
        circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(5), Inches(4), Inches(4))
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.primary, 0.1))
        circle1.line.fill.background()
        
        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(4), Inches(4))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = hex_to_rgb(self.accent)
        circle2.line.fill.background()
        
        # Title in white
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title.split('.')[0].strip()[:50]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.title_font
        
        # Stats cards
        stats = body_texts[:4]
        card_width = 2.8
        total_width = len(stats) * card_width + (len(stats) - 1) * 0.4
        start_x = (13.333 - total_width) / 2
        
        for i, stat_text in enumerate(stats):
            x = start_x + i * (card_width + 0.4)
            
            # White card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2.2),
                Inches(card_width), Inches(3.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.fill.background()
            
            # Parse stat
            parts = stat_text.split(':', 1) if ':' in stat_text else [stat_text[:20], stat_text[20:]]
            value = parts[0].strip()
            label = parts[1].strip() if len(parts) > 1 else ''
            
            # Big number/value
            value_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.6),
                Inches(card_width - 0.2), Inches(1.2)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value[:15]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.primary)
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.title_font
            
            # Label
            if label:
                label_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(3.9),
                    Inches(card_width - 0.2), Inches(1.5)
                )
                tf = label_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = label[:60]
                p.font.size = Pt(11)
                p.font.color.rgb = hex_to_rgb(self.text_muted)
                p.alignment = PP_ALIGN.CENTER
                p.font.name = self.body_font
    
    # ==================== STUNNING CLOSING SLIDE ====================
    
    def _create_stunning_closing_slide(self, title: str, body_texts: List[str]):
        """Create a memorable closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, self.primary)
        
        # Large decorative shapes
        shape1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9), Inches(4),
            Inches(6), Inches(6)
        )
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = hex_to_rgb(self.accent)
        shape1.line.fill.background()
        
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-2), Inches(-3),
            Inches(8), Inches(8)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.primary, 0.15))
        shape2.line.fill.background()
        
        # Main message
        message = title if title else "Thank You"
        if "terima" in message.lower() or "thank" in message.lower():
            message = "Thank You"
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = message
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.title_font
        
        # Tagline
        tagline = body_texts[0] if body_texts else "Questions?"
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline[:80]
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.body_font
        
        # Decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(5.2),
            Inches(2.333), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.accent)
        line.line.fill.background()
    
    # ==================== UTILITIES ====================
    
    def _set_slide_background(self, slide, color: str):
        """Set slide background color"""
        if color.startswith('#'):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(color)
    
    def _is_dark_color(self, color: str) -> bool:
        """Check if a color is dark"""
        if color.startswith('#'):
            hex_color = color.lstrip('#')
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
                return luminance < 0.5
        return False


# Keep old class name for backwards compatibility
PPTXExporter = WorldClassExporter


def export_presentation(style: Dict, slides_data: List[Dict], output_path: str) -> str:
    """Convenience function to export a presentation"""
    exporter = WorldClassExporter(style, slides_data)
    return exporter.export(output_path)
