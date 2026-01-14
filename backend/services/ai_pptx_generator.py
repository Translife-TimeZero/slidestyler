"""
AI-Powered PPTX Generator with Seedream-4 Image Integration
Creates stunning presentations where AI images and text blend perfectly
"""

import os
import io
import json
import httpx
import asyncio
import base64
import tempfile
from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    return RGBColor(0, 0, 0)


# Prompt for AI to generate design with image integration
DESIGN_WITH_IMAGE_PROMPT = """You are a world-class presentation designer. Analyze this slide content and create a design where text and visuals blend perfectly.

SLIDE CONTENT:
{content}

SLIDE POSITION: {position} of {total} slides
CONTEXT: {context}

Generate a JSON response:
{{
    "layout": "hero_image|split_image|corner_image|full_background|cards|grid",
    "image_prompt": "describe an abstract, professional image that represents this content (e.g., 'abstract blue waves representing flow and progress', 'geometric shapes in dark blue suggesting structure')",
    "image_position": "background|left|right|top-right|bottom-left",
    "image_opacity": 0.3-1.0,
    "text_area": {{
        "position": "left|right|center|bottom",
        "has_overlay": true/false,
        "overlay_color": "#hex with 80% opacity"
    }},
    "colors": {{
        "primary": "#hex (main brand color)",
        "accent": "#hex (highlight color)", 
        "text": "#hex (must contrast with background)",
        "text_secondary": "#hex (subtitles, muted)"
    }},
    "title": {{
        "text": "clean, shortened title",
        "size": 32-54,
        "position": "specify x,y as percentages"
    }},
    "mood": "professional|dynamic|calm|bold|innovative",
    "blend_style": "overlay|side-by-side|floating-cards|gradient-fade"
}}

CRITICAL RULES:
- Text MUST be readable - use overlays or position text away from busy image areas
- Image should enhance message, not distract
- First slide = impactful hero with dramatic image
- Last slide = memorable closing
- Colors must work together harmoniously

Respond ONLY with valid JSON."""


class AIPPTXGenerator:
    """AI-powered PPTX generator with Seedream-4 image integration"""
    
    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key or os.environ.get("REPLICATE_API_TOKEN")
        self.llm_url = "https://api.replicate.com/v1/models/meta/meta-llama-3-70b-instruct/predictions"
        self.image_url = "https://api.replicate.com/v1/models/bytedance/seedream-3.0/predictions"
        self.presentation_context = ""
        self.generated_images = {}  # Cache generated images
        
    async def generate_presentation(
        self, 
        slides_data: List[Dict],
        output_path: str,
        generate_images: bool = True
    ) -> str:
        """Generate presentation with AI-designed slides and images"""
        
        # Create presentation (16:9 widescreen)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Build context from all slides
        all_content = []
        for slide in slides_data:
            texts = self._extract_texts(slide.get('original_content', []))
            all_content.extend(texts[:3])
        self.presentation_context = " | ".join(all_content[:10])[:500]
        
        # Determine overall theme/mood for consistency
        overall_mood = await self._analyze_overall_mood(self.presentation_context)
        
        # Generate each slide
        total = len(slides_data)
        for i, slide_data in enumerate(slides_data):
            print(f"[AI Generator] Creating slide {i+1}/{total}...")
            
            # Get AI design with image instructions
            design = await self._get_ai_design_with_image(slide_data, i, total, overall_mood)
            
            # Generate image if needed
            image_path = None
            if generate_images and design.get('image_prompt'):
                image_path = await self._generate_slide_image(design, i)
            
            # Create the slide with perfect text-image blend
            self._create_blended_slide(prs, slide_data, design, image_path, i, total)
        
        prs.save(output_path)
        return output_path
    
    def _extract_texts(self, original_content: List) -> List[str]:
        """Extract text items from content"""
        texts = []
        skip_types = ['sldNum', 'ftr', 'dt', 'hdr']
        
        for item in original_content:
            if isinstance(item, str):
                text = item.strip()
            else:
                if item.get('type') in skip_types:
                    continue
                text = item.get('text', '').strip()
            
            if text and not (text.isdigit() and len(text) <= 3):
                texts.append(text)
        
        return texts
    
    async def _analyze_overall_mood(self, context: str) -> Dict:
        """Analyze the overall presentation mood for consistency"""
        # Default mood palette
        return {
            "primary": "#1e3a5f",
            "accent": "#00d4aa",
            "mood": "professional",
            "image_style": "abstract geometric shapes, professional, modern, clean lines"
        }
    
    async def _get_ai_design_with_image(
        self, 
        slide_data: Dict, 
        index: int, 
        total: int,
        overall_mood: Dict
    ) -> Dict:
        """Get AI design instructions including image generation prompt"""
        texts = self._extract_texts(slide_data.get('original_content', []))
        content_summary = "\n".join([f"- {t[:100]}" for t in texts[:8]])
        
        position = "first (title slide)" if index == 0 else \
                   "last (closing slide)" if index == total - 1 else \
                   f"middle ({index + 1})"
        
        # Try AI, fallback to smart defaults
        try:
            if self.api_key:
                design = await self._call_llm(content_summary, position, total)
                if design:
                    return design
        except Exception as e:
            print(f"[AI Generator] LLM call failed: {e}")
        
        # Smart fallback with image prompts
        return self._get_smart_design(index, total, texts, overall_mood)
    
    async def _call_llm(self, content: str, position: str, total: int) -> Optional[Dict]:
        """Call LLM for design instructions"""
        prompt = DESIGN_WITH_IMAGE_PROMPT.format(
            content=content,
            position=position,
            total=total,
            context=self.presentation_context[:300]
        )
        
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                self.llm_url,
                headers={
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                },
                json={"input": {"prompt": prompt, "max_tokens": 1000, "temperature": 0.7}}
            )
            
            if response.status_code != 201:
                raise Exception(f"API error: {response.status_code}")
            
            result = response.json()
            prediction_url = result.get('urls', {}).get('get')
            
            for _ in range(30):
                await asyncio.sleep(1)
                poll = await client.get(prediction_url, headers={"Authorization": f"Bearer {self.api_key}"})
                poll_result = poll.json()
                
                if poll_result.get('status') == 'succeeded':
                    output = poll_result.get('output', '')
                    if isinstance(output, list):
                        output = ''.join(output)
                    return self._parse_json(output)
                elif poll_result.get('status') == 'failed':
                    raise Exception("Prediction failed")
            
            raise Exception("Timeout")
    
    def _parse_json(self, text: str) -> Dict:
        """Parse JSON from AI response"""
        text = text.strip()
        if '```json' in text:
            start = text.find('```json') + 7
            end = text.find('```', start)
            text = text[start:end].strip()
        elif '```' in text:
            start = text.find('```') + 3
            end = text.find('```', start)
            text = text[start:end].strip()
        
        start = text.find('{')
        end = text.rfind('}') + 1
        if start != -1 and end > start:
            return json.loads(text[start:end])
        raise ValueError("No JSON found")
    
    def _get_smart_design(self, index: int, total: int, texts: List[str], mood: Dict) -> Dict:
        """Generate smart design with image prompts based on content"""
        is_first = index == 0
        is_last = index == total - 1
        
        # Analyze content for image prompt
        content_text = " ".join(texts[:3]).lower()
        
        # Generate contextual image prompts
        image_prompts = {
            "learning": "abstract flowing lines representing knowledge transfer, blue and teal gradient, professional",
            "error": "geometric pattern showing correction and improvement, structured shapes, warm colors",
            "process": "interconnected nodes and pathways, systematic flow diagram style, modern blue",
            "leadership": "ascending geometric shapes suggesting growth and direction, bold colors",
            "team": "overlapping circles representing collaboration, harmonious colors",
            "change": "dynamic arrows and transformation shapes, gradient from old to new",
            "analysis": "data visualization abstract, charts and graphs stylized, professional blue",
            "default": "abstract professional background, subtle geometric patterns, corporate blue gradient"
        }
        
        # Choose image prompt based on content
        image_prompt = image_prompts["default"]
        for keyword, prompt in image_prompts.items():
            if keyword in content_text:
                image_prompt = prompt
                break
        
        # Color palettes that work well with images
        palettes = [
            {"bg": "#0f172a", "primary": "#3b82f6", "accent": "#22d3ee", "text": "#ffffff", "overlay": "#0f172aCC"},
            {"bg": "#1e1b4b", "primary": "#8b5cf6", "accent": "#f472b6", "text": "#ffffff", "overlay": "#1e1b4bCC"},
            {"bg": "#14532d", "primary": "#22c55e", "accent": "#fbbf24", "text": "#ffffff", "overlay": "#14532dCC"},
            {"bg": "#7c2d12", "primary": "#f97316", "accent": "#fcd34d", "text": "#ffffff", "overlay": "#7c2d12CC"},
            {"bg": "#1e3a5f", "primary": "#0ea5e9", "accent": "#2dd4bf", "text": "#ffffff", "overlay": "#1e3a5fCC"},
        ]
        palette = palettes[index % len(palettes)]
        
        if is_first:
            return {
                "layout": "hero_image",
                "image_prompt": f"dramatic {image_prompt}, wide cinematic composition, dark overlay ready",
                "image_position": "background",
                "image_opacity": 0.4,
                "text_area": {"position": "center", "has_overlay": True, "overlay_color": palette["overlay"]},
                "colors": {"primary": palette["primary"], "accent": palette["accent"], "text": "#ffffff", "text_secondary": "#ffffffCC"},
                "title": {"text": texts[0][:50] if texts else "Presentation", "size": 52, "position": "center"},
                "mood": "bold",
                "blend_style": "gradient-fade"
            }
        elif is_last:
            return {
                "layout": "hero_image",
                "image_prompt": f"inspiring {image_prompt}, uplifting composition, space for text in center",
                "image_position": "background",
                "image_opacity": 0.3,
                "text_area": {"position": "center", "has_overlay": True, "overlay_color": palette["overlay"]},
                "colors": {"primary": palette["primary"], "accent": palette["accent"], "text": "#ffffff", "text_secondary": "#ffffffCC"},
                "title": {"text": "Thank You", "size": 54, "position": "center"},
                "mood": "professional",
                "blend_style": "gradient-fade"
            }
        elif len(texts) > 5:
            return {
                "layout": "corner_image",
                "image_prompt": f"small accent {image_prompt}, corner composition, minimal",
                "image_position": "top-right",
                "image_opacity": 0.8,
                "text_area": {"position": "left", "has_overlay": False, "overlay_color": None},
                "colors": {"primary": palette["primary"], "accent": palette["accent"], "text": "#ffffff", "text_secondary": "#ffffffB3"},
                "title": {"text": texts[0][:45] if texts else "", "size": 32, "position": "top-left"},
                "mood": "professional",
                "blend_style": "side-by-side"
            }
        else:
            return {
                "layout": "split_image",
                "image_prompt": f"{image_prompt}, vertical composition, clean edges",
                "image_position": "right",
                "image_opacity": 1.0,
                "text_area": {"position": "left", "has_overlay": False, "overlay_color": None},
                "colors": {"primary": palette["primary"], "accent": palette["accent"], "text": "#ffffff", "text_secondary": "#ffffffB3"},
                "title": {"text": texts[0][:45] if texts else "", "size": 32, "position": "top-left"},
                "mood": "professional",
                "blend_style": "side-by-side"
            }
    
    async def _generate_slide_image(self, design: Dict, slide_index: int) -> Optional[str]:
        """Generate image using Seedream-3"""
        if not self.api_key:
            return None
        
        image_prompt = design.get('image_prompt', '')
        if not image_prompt:
            return None
        
        # Enhance prompt for better results
        enhanced_prompt = f"{image_prompt}, high quality, 4k, professional presentation background, no text, abstract"
        
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(
                    self.image_url,
                    headers={
                        "Authorization": f"Bearer {self.api_key}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "input": {
                            "prompt": enhanced_prompt,
                            "num_outputs": 1,
                            "aspect_ratio": "16:9",
                            "output_format": "png"
                        }
                    }
                )
                
                if response.status_code != 201:
                    print(f"[Image Gen] API error: {response.status_code}")
                    return None
                
                result = response.json()
                prediction_url = result.get('urls', {}).get('get')
                
                # Poll for result
                for _ in range(60):  # Wait up to 60 seconds
                    await asyncio.sleep(2)
                    poll = await client.get(
                        prediction_url,
                        headers={"Authorization": f"Bearer {self.api_key}"}
                    )
                    poll_result = poll.json()
                    
                    if poll_result.get('status') == 'succeeded':
                        output = poll_result.get('output')
                        if output:
                            image_url = output[0] if isinstance(output, list) else output
                            # Download and save image
                            img_response = await client.get(image_url)
                            if img_response.status_code == 200:
                                temp_path = tempfile.mktemp(suffix='.png')
                                with open(temp_path, 'wb') as f:
                                    f.write(img_response.content)
                                print(f"[Image Gen] Generated image for slide {slide_index + 1}")
                                return temp_path
                    elif poll_result.get('status') == 'failed':
                        print(f"[Image Gen] Generation failed")
                        return None
                
                print(f"[Image Gen] Timeout for slide {slide_index + 1}")
                return None
                
        except Exception as e:
            print(f"[Image Gen] Error: {e}")
            return None
    
    def _create_blended_slide(
        self, 
        prs: Presentation, 
        slide_data: Dict, 
        design: Dict, 
        image_path: Optional[str],
        index: int,
        total: int
    ):
        """Create a slide with perfectly blended image and text"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        texts = self._extract_texts(slide_data.get('original_content', []))
        
        colors = design.get('colors', {})
        bg_color = colors.get('primary', '#1e3a5f')
        text_color = colors.get('text', '#ffffff')
        accent_color = colors.get('accent', '#22d3ee')
        
        layout = design.get('layout', 'cards')
        image_position = design.get('image_position', 'right')
        text_area = design.get('text_area', {})
        
        # Step 1: Set background color
        self._set_background(slide, bg_color)
        
        # Step 2: Add image if available
        if image_path and os.path.exists(image_path):
            self._add_blended_image(slide, image_path, design)
        else:
            # Add decorative shapes as fallback
            self._add_decorative_shapes(slide, colors, layout)
        
        # Step 3: Add overlay for text readability if needed
        if text_area.get('has_overlay') and text_area.get('overlay_color'):
            self._add_text_overlay(slide, design)
        
        # Step 4: Add text content in designated area
        title_config = design.get('title', {})
        title_text = title_config.get('text', texts[0] if texts else '')
        
        self._add_blended_text(slide, title_text, texts[1:], design, colors)
        
        # Step 5: Add accent elements
        self._add_accent_elements(slide, colors, layout, index, total)
    
    def _set_background(self, slide, color: str):
        """Set slide background"""
        if color.startswith('#'):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(color)
    
    def _add_blended_image(self, slide, image_path: str, design: Dict):
        """Add image with proper positioning for text blend"""
        position = design.get('image_position', 'background')
        opacity = design.get('image_opacity', 0.5)
        
        if position == 'background':
            # Full background image
            slide.shapes.add_picture(
                image_path,
                Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5)
            )
        elif position == 'right':
            # Right side (60% width)
            slide.shapes.add_picture(
                image_path,
                Inches(5.5), Inches(0),
                width=Inches(7.833), height=Inches(7.5)
            )
        elif position == 'left':
            # Left side
            slide.shapes.add_picture(
                image_path,
                Inches(0), Inches(0),
                width=Inches(6.5), height=Inches(7.5)
            )
        elif position == 'top-right':
            # Corner accent
            slide.shapes.add_picture(
                image_path,
                Inches(8), Inches(0),
                width=Inches(5.333), height=Inches(4)
            )
        elif position == 'bottom-left':
            slide.shapes.add_picture(
                image_path,
                Inches(0), Inches(4),
                width=Inches(5), height=Inches(3.5)
            )
    
    def _add_decorative_shapes(self, slide, colors: Dict, layout: str):
        """Add decorative shapes when no image is available"""
        accent = colors.get('accent', '#22d3ee')
        primary = colors.get('primary', '#1e3a5f')
        
        if layout in ['hero_image', 'full_background']:
            # Large accent circle (top-right)
            circle1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(9), Inches(-2),
                Inches(6), Inches(6)
            )
            circle1.fill.solid()
            circle1.fill.fore_color.rgb = hex_to_rgb(accent)
            circle1.fill.fore_color.brightness = 0.3
            circle1.line.fill.background()
            
            # Small accent (bottom-left)
            circle2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(-1), Inches(5),
                Inches(4), Inches(4)
            )
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = hex_to_rgb(accent)
            circle2.fill.fore_color.brightness = 0.2
            circle2.line.fill.background()
        
        elif layout == 'split_image':
            # Right panel
            panel = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(6.5), Inches(0),
                Inches(6.833), Inches(7.5)
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = hex_to_rgb(accent)
            panel.fill.fore_color.brightness = 0.4
            panel.line.fill.background()
            
            # Decorative circles on panel
            for i, (x, y, s) in enumerate([(8, 1, 2), (10, 4, 1.5), (7, 5, 1)]):
                c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(255, 255, 255)
                c.fill.fore_color.brightness = 0.7 + i * 0.1
                c.line.fill.background()
        
        elif layout == 'corner_image':
            # Corner accent
            corner = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(10), Inches(-1),
                Inches(4.5), Inches(4.5)
            )
            corner.fill.solid()
            corner.fill.fore_color.rgb = hex_to_rgb(accent)
            corner.fill.fore_color.brightness = 0.3
            corner.line.fill.background()
    
    def _add_text_overlay(self, slide, design: Dict):
        """Add semi-transparent overlay for text readability"""
        text_area = design.get('text_area', {})
        position = text_area.get('position', 'center')
        
        if position == 'center':
            # Center gradient overlay
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(1.5),
                Inches(13.333), Inches(4.5)
            )
        elif position == 'left':
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(6), Inches(7.5)
            )
        elif position == 'bottom':
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(4),
                Inches(13.333), Inches(3.5)
            )
        else:
            return
        
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.fore_color.brightness = -0.5
        overlay.line.fill.background()
    
    def _add_blended_text(self, slide, title: str, body_texts: List[str], design: Dict, colors: Dict):
        """Add text positioned to blend with image"""
        text_area = design.get('text_area', {})
        position = text_area.get('position', 'left')
        layout = design.get('layout', 'cards')
        
        text_color = colors.get('text', '#ffffff')
        accent_color = colors.get('accent', '#22d3ee')
        title_size = design.get('title', {}).get('size', 32)
        
        # Determine text positions based on layout
        if layout in ['hero_image', 'full_background'] or position == 'center':
            title_x, title_y, title_w = 0.8, 2.5, 11.7
            body_x, body_y, body_w = 0.8, 4.0, 11.7
            center_align = True
        elif position == 'left' or layout == 'split_image':
            title_x, title_y, title_w = 0.6, 0.8, 5.5
            body_x, body_y, body_w = 0.6, 1.8, 5.5
            center_align = False
        else:
            title_x, title_y, title_w = 0.6, 0.6, 7
            body_x, body_y, body_w = 0.6, 1.6, 7
            center_align = False
        
        # Add title
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(title_x), Inches(title_y),
                Inches(title_w), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title[:70]
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(text_color)
            if center_align:
                p.alignment = PP_ALIGN.CENTER
        
        # Add body content based on layout
        if layout in ['hero_image', 'full_background']:
            # Simple subtitle for hero slides
            if body_texts:
                sub_box = slide.shapes.add_textbox(
                    Inches(body_x), Inches(body_y + 0.3),
                    Inches(body_w), Inches(1)
                )
                tf = sub_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = body_texts[0][:100]
                p.font.size = Pt(20)
                p.font.color.rgb = hex_to_rgb(colors.get('text_secondary', '#ffffffCC'))
                if center_align:
                    p.alignment = PP_ALIGN.CENTER
        else:
            # Cards or list for content slides
            self._add_content_cards(slide, body_texts[:6], body_x, body_y, body_w, colors)
    
    def _add_content_cards(
        self, 
        slide, 
        items: List[str], 
        x: float, 
        y: float, 
        width: float,
        colors: Dict
    ):
        """Add content as elegant cards"""
        accent = colors.get('accent', '#22d3ee')
        card_height = 0.85
        
        for i, text in enumerate(items):
            card_y = y + i * (card_height + 0.12)
            if card_y > 6.2:
                break
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(card_y),
                Inches(width), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.fill.fore_color.brightness = 0.85
            card.line.fill.background()
            
            # Left accent bar
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(card_y),
                Inches(0.08), Inches(card_height)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = hex_to_rgb(accent)
            bar.line.fill.background()
            
            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.2), Inches(card_y + 0.18),
                Inches(0.5), Inches(0.5)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex_to_rgb(accent)
            badge.line.fill.background()
            
            # Number text
            num_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(card_y + 0.22),
                Inches(0.5), Inches(0.45)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Card text
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.85), Inches(card_y + 0.18),
                Inches(width - 1.1), Inches(card_height - 0.35)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text[:110]
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(30, 30, 50)
    
    def _add_accent_elements(self, slide, colors: Dict, layout: str, index: int, total: int):
        """Add finishing accent elements"""
        accent = colors.get('accent', '#22d3ee')
        
        # Top accent line (except for hero slides)
        if layout not in ['hero_image', 'full_background']:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(0.08)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = hex_to_rgb(accent)
            line.line.fill.background()
        
        # Bottom line for closing
        if index == total - 1:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.5), Inches(5.5),
                Inches(2.333), Inches(0.06)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = hex_to_rgb(accent)
            line.line.fill.background()


async def generate_ai_presentation(
    slides_data: List[Dict], 
    output_path: str, 
    api_key: str = None,
    generate_images: bool = True
) -> str:
    """Convenience function to generate AI-designed presentation with images"""
    generator = AIPPTXGenerator(api_key)
    return await generator.generate_presentation(slides_data, output_path, generate_images)
